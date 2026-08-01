"""
Quality check implementations for generated reports.

Validates PPTX and Excel outputs to ensure they meet requirements
before delivery.
"""

import logging
from io import BytesIO
from typing import Literal

from pydantic import BaseModel, Field
from pptx import Presentation as PptxPresentation
from openpyxl import load_workbook

from src.modules.ai_insight.models import PresentationPlan

logger = logging.getLogger(__name__)


class CheckItem(BaseModel):
    """Single quality check result."""
    name: str
    status: Literal["pass", "fail", "warn"]
    message: str


class QualityCheckResult(BaseModel):
    """Aggregate quality check results."""
    passed: bool
    total_checks: int = 0
    passed_checks: int = 0
    failed_checks: int = 0
    warned_checks: int = 0
    checks: list[CheckItem] = Field(default_factory=list)


def check_slide_count(pptx_bytes: bytes, expected: int = 16) -> CheckItem:
    """Verify the PPTX has exactly the expected number of slides."""
    try:
        prs = PptxPresentation(BytesIO(pptx_bytes))
        actual = len(prs.slides)
        if actual == expected:
            return CheckItem(
                name="slide_count",
                status="pass",
                message=f"Slide count OK: {actual} slides",
            )
        else:
            return CheckItem(
                name="slide_count",
                status="fail",
                message=f"Expected {expected} slides, got {actual}",
            )
    except Exception as e:
        return CheckItem(
            name="slide_count",
            status="fail",
            message=f"Failed to open PPTX: {e}",
        )


def check_charts_present(pptx_bytes: bytes, plan: PresentationPlan) -> list[CheckItem]:
    """Verify that slides with expected charts actually contain chart shapes."""
    results = []
    try:
        prs = PptxPresentation(BytesIO(pptx_bytes))
        slides = list(prs.slides)

        for slide_plan in plan.slides:
            if not slide_plan.chart:
                continue

            idx = slide_plan.page_number - 1
            if idx >= len(slides):
                results.append(CheckItem(
                    name=f"chart_page_{slide_plan.page_number}",
                    status="fail",
                    message=f"Page {slide_plan.page_number} does not exist",
                ))
                continue

            slide = slides[idx]
            has_chart = any(shape.has_chart for shape in slide.shapes)

            if has_chart:
                results.append(CheckItem(
                    name=f"chart_page_{slide_plan.page_number}",
                    status="pass",
                    message=f"Page {slide_plan.page_number} has native chart",
                ))
            else:
                results.append(CheckItem(
                    name=f"chart_page_{slide_plan.page_number}",
                    status="fail",
                    message=f"Page {slide_plan.page_number} missing expected chart",
                ))

    except Exception as e:
        results.append(CheckItem(
            name="chart_presence",
            status="fail",
            message=f"Failed to check charts: {e}",
        ))

    return results


def check_no_image_charts(pptx_bytes: bytes) -> CheckItem:
    """
    Verify no slides use image-based charts (competition requirement).
    Charts must be native vector objects, not pasted images.
    """
    try:
        prs = PptxPresentation(BytesIO(pptx_bytes))
        image_count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    # Check if it looks like a chart image
                    if hasattr(shape, "image"):
                        image_count += 1

        if image_count == 0:
            return CheckItem(
                name="no_image_charts",
                status="pass",
                message="No image-based charts detected (all native)",
            )
        else:
            return CheckItem(
                name="no_image_charts",
                status="warn",
                message=f"Found {image_count} images — verify none are chart screenshots",
            )
    except Exception as e:
        return CheckItem(
            name="no_image_charts",
            status="warn",
            message=f"Could not verify image charts: {e}",
        )


def check_text_completeness(pptx_bytes: bytes) -> CheckItem:
    """Check that no slides have empty text placeholders."""
    try:
        prs = PptxPresentation(BytesIO(pptx_bytes))
        empty_slides = []

        for i, slide in enumerate(prs.slides, 1):
            has_text = False
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    has_text = True
                    break
            if not has_text:
                empty_slides.append(i)

        if not empty_slides:
            return CheckItem(
                name="text_completeness",
                status="pass",
                message="All slides have text content",
            )
        else:
            return CheckItem(
                name="text_completeness",
                status="fail",
                message=f"Empty slides detected: {empty_slides}",
            )
    except Exception as e:
        return CheckItem(
            name="text_completeness",
            status="fail",
            message=f"Failed text check: {e}",
        )


def check_excel_sheets(xlsx_bytes: bytes) -> CheckItem:
    """Verify Excel has expected worksheet structure."""
    try:
        wb = load_workbook(BytesIO(xlsx_bytes), read_only=True)
        sheet_names = wb.sheetnames

        if len(sheet_names) >= 2:
            return CheckItem(
                name="excel_sheets",
                status="pass",
                message=f"Excel has {len(sheet_names)} sheets: {', '.join(sheet_names)}",
            )
        else:
            return CheckItem(
                name="excel_sheets",
                status="warn",
                message=f"Excel only has {len(sheet_names)} sheet(s)",
            )
    except Exception as e:
        return CheckItem(
            name="excel_sheets",
            status="fail",
            message=f"Failed to open Excel: {e}",
        )


def check_excel_charts(xlsx_bytes: bytes) -> CheckItem:
    """Verify Excel contains native chart objects."""
    try:
        wb = load_workbook(BytesIO(xlsx_bytes))
        chart_count = 0

        for ws in wb.worksheets:
            chart_count += len(ws._charts)

        if chart_count > 0:
            return CheckItem(
                name="excel_charts",
                status="pass",
                message=f"Excel has {chart_count} native chart(s)",
            )
        else:
            return CheckItem(
                name="excel_charts",
                status="warn",
                message="No native charts found in Excel",
            )
    except Exception as e:
        return CheckItem(
            name="excel_charts",
            status="fail",
            message=f"Failed chart check: {e}",
        )


def check_file_size(pptx_bytes: bytes, xlsx_bytes: bytes, max_mb: int = 50) -> CheckItem:
    """Verify files are not unreasonably large."""
    pptx_mb = len(pptx_bytes) / (1024 * 1024)
    xlsx_mb = len(xlsx_bytes) / (1024 * 1024)
    total_mb = pptx_mb + xlsx_mb

    if total_mb < max_mb:
        return CheckItem(
            name="file_size",
            status="pass",
            message=f"Total size OK: PPTX={pptx_mb:.1f}MB, Excel={xlsx_mb:.1f}MB",
        )
    else:
        return CheckItem(
            name="file_size",
            status="warn",
            message=f"Large files: PPTX={pptx_mb:.1f}MB, Excel={xlsx_mb:.1f}MB (total {total_mb:.1f}MB)",
        )


def run_all_checks(
    pptx_bytes: bytes,
    xlsx_bytes: bytes,
    plan: PresentationPlan,
) -> QualityCheckResult:
    """
    Run all quality checks and return aggregate result.

    Args:
        pptx_bytes: Generated PPTX file bytes
        xlsx_bytes: Generated Excel file bytes
        plan: Original presentation plan for comparison

    Returns:
        QualityCheckResult with pass/fail determination
    """
    checks: list[CheckItem] = []

    # PPTX checks
    checks.append(check_slide_count(pptx_bytes))
    checks.extend(check_charts_present(pptx_bytes, plan))
    checks.append(check_no_image_charts(pptx_bytes))
    checks.append(check_text_completeness(pptx_bytes))

    # Excel checks
    checks.append(check_excel_sheets(xlsx_bytes))
    checks.append(check_excel_charts(xlsx_bytes))

    # Size check
    checks.append(check_file_size(pptx_bytes, xlsx_bytes))

    # Aggregate
    passed_count = sum(1 for c in checks if c.status == "pass")
    failed_count = sum(1 for c in checks if c.status == "fail")
    warned_count = sum(1 for c in checks if c.status == "warn")

    result = QualityCheckResult(
        passed=failed_count == 0,
        total_checks=len(checks),
        passed_checks=passed_count,
        failed_checks=failed_count,
        warned_checks=warned_count,
        checks=checks,
    )

    logger.info(
        f"Quality checks: {passed_count} pass, {failed_count} fail, "
        f"{warned_count} warn — {'PASSED' if result.passed else 'FAILED'}"
    )

    return result
