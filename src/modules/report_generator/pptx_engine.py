"""
Native PowerPoint generation engine with smart layout selection.

Optimized for TS Holdings template:
- Safe content area: avoid bottom red gradient + logo zone (below Y=5.8")
- Title area: Y 0.3" to 1.2"
- Content area: Y 1.4" to 5.6"
- Speaker notes: AI analysis reasoning added to every content slide
"""

import logging
import copy
from io import BytesIO

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from src.modules.ai_insight.models import PresentationPlan, SlideContent, ChartSpec

logger = logging.getLogger(__name__)

PPTX_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "waterfall": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "heatmap": XL_CHART_TYPE.COLUMN_CLUSTERED,
}

# Safe content zones (avoid template decorations)
# Template is 13.333" x 7.5" (widescreen 16:9)
TITLE_TOP = Inches(0.4)
TITLE_HEIGHT = Inches(0.9)
CONTENT_TOP = Inches(1.5)
CONTENT_BOTTOM = Inches(5.6)  # Stop before bottom decoration
CONTENT_HEIGHT = Inches(4.1)  # 5.6 - 1.5
LEFT_MARGIN = Inches(0.6)
RIGHT_MARGIN = Inches(0.6)
SLIDE_WIDTH = Inches(13.333)
FULL_CONTENT_WIDTH = Inches(12.1)  # 13.333 - 0.6 - 0.6


class PptxGenerator:
    """Generates PPTX with native charts, smart layout selection, and speaker notes.
    
    Template strategy (v3 - template-first):
    - If a template is provided, keep its slides intact as backgrounds
    - Duplicate template slides to reach 16 pages
    - Clear text placeholders on each slide, then fill with new AI content
    - Add speaker notes last
    - This preserves all template decorations (images, gradients, logos)
    """

    def __init__(self, plan: PresentationPlan, template_bytes: bytes | None = None):
        self._plan = plan
        self._has_template = template_bytes is not None

        if template_bytes:
            self._prs = Presentation(BytesIO(template_bytes))
            num_template_slides = len(self._prs.slides)
            logger.info(f"Template loaded with {num_template_slides} slides")

            # Categorize template slides by their visual purpose
            self._template_slide_categories = self._categorize_template_slides()

            # Build 16 slides from template: duplicate existing slides as needed
            self._prepare_16_slides_from_template(num_template_slides)
        else:
            self._prs = Presentation()
            self._prs.slide_width = Inches(13.333)
            self._prs.slide_height = Inches(7.5)
            self._template_slide_categories = {}

        if not self._has_template:
            self._layout_map = self._analyze_layouts()
            logger.info(f"Layout mapping: {self._layout_map}")

    def _categorize_template_slides(self):
        """Categorize each template slide as 'cover', 'content', or 'divider'."""
        categories = {}
        for idx, slide in enumerate(self._prs.slides):
            # Count shapes to estimate slide type
            shape_count = len(slide.shapes)
            has_title_ph = slide.shapes.title is not None
            
            # Heuristic: first slide is cover, last might be ending
            # Slides with fewer shapes are likely content templates
            if idx == 0:
                categories[idx] = 'cover'
            elif shape_count > 5:
                categories[idx] = 'decorated'  # heavily decorated (divider/cover)
            else:
                categories[idx] = 'content'  # cleaner layout for content
        
        logger.info(f"Template slide categories: {categories}")
        return categories

    def _prepare_16_slides_from_template(self, num_template_slides: int):
        """
        Ensure we have exactly 16 slides by duplicating template slides.
        Strategy:
        - Keep slide 0 as cover (page 1)
        - Use the cleanest slides for content pages
        - Duplicate as needed to reach 16
        """
        target = 16
        
        if num_template_slides >= target:
            # Template has enough slides, remove extras
            while len(self._prs.slides) > target:
                rId = self._prs.slides._sldIdLst[-1].rId
                self._prs.part.drop_rel(rId)
                del self._prs.slides._sldIdLst[-1]
            logger.info(f"Template trimmed to {target} slides")
            return

        # Need to duplicate slides to reach 16
        # Find the best "content" slide to duplicate (least decorated)
        best_content_idx = self._find_best_content_slide()
        
        # Duplicate slides until we have 16
        slides_to_add = target - num_template_slides
        logger.info(f"Duplicating slide [{best_content_idx}] {slides_to_add} times to reach {target}")
        
        for _ in range(slides_to_add):
            self._duplicate_slide(best_content_idx)

    def _find_best_content_slide(self):
        """Find the template slide best suited for content (least decorated, not cover)."""
        best_idx = min(1, len(self._prs.slides) - 1)  # Default: second slide
        min_shapes = 999
        
        for idx, slide in enumerate(self._prs.slides):
            if idx == 0:
                continue  # Skip cover
            shape_count = len(slide.shapes)
            if shape_count < min_shapes:
                min_shapes = shape_count
                best_idx = idx
        
        return best_idx

    def _duplicate_slide(self, src_idx: int):
        """Duplicate a slide (preserving all shapes, images, and relationships)."""
        try:
            src_slide = self._prs.slides[src_idx]
            # Use the same layout as the source slide
            layout = src_slide.slide_layout
            new_slide = self._prs.slides.add_slide(layout)
            
            # Copy all shapes from source to new slide
            src_sp_tree = src_slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
            dst_sp_tree = new_slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
            
            # Remove default placeholder shapes from new slide
            for child in list(dst_sp_tree):
                if child.tag not in [qn('p:nvGrpSpPr'), qn('p:grpSpPr')]:
                    dst_sp_tree.remove(child)
            
            # Copy shapes from source (including images via relationships)
            for child in src_sp_tree:
                if child.tag not in [qn('p:nvGrpSpPr'), qn('p:grpSpPr')]:
                    dst_sp_tree.append(copy.deepcopy(child))
            
            # Copy background if exists
            src_cSld = src_slide._element.find(qn('p:cSld'))
            dst_cSld = new_slide._element.find(qn('p:cSld'))
            src_bg = src_cSld.find(qn('p:bg'))
            if src_bg is not None:
                dst_bg = dst_cSld.find(qn('p:bg'))
                if dst_bg is not None:
                    dst_cSld.remove(dst_bg)
                dst_cSld.insert(0, copy.deepcopy(src_bg))

        except Exception as e:
            logger.warning(f"Failed to duplicate slide {src_idx}: {e}")
            # Fallback: just add a blank slide with the layout
            layout = self._prs.slides[src_idx].slide_layout
            self._prs.slides.add_slide(layout)

    def _analyze_layouts(self):
        layout_map = {}
        layouts = self._prs.slide_layouts
        for idx, layout in enumerate(layouts):
            name = (layout.name or "").lower()
            logger.info(f"  [{idx}] '{layout.name}'")
            if any(kw in name for kw in ["標題投影片", "title slide"]):
                if "title_cover" not in layout_map:
                    layout_map["title_cover"] = idx
            elif any(kw in name for kw in ["區段標頭", "section header"]):
                if "section" not in layout_map:
                    layout_map["section"] = idx
            elif any(kw in name for kw in ["僅標題", "title only"]):
                if "title_only" not in layout_map:
                    layout_map["title_only"] = idx
            elif any(kw in name for kw in ["空白", "blank"]):
                if "blank" not in layout_map:
                    layout_map["blank"] = idx
            elif any(kw in name for kw in ["標題及內容", "title and content"]):
                if "title_content" not in layout_map:
                    layout_map["title_content"] = idx
            elif any(kw in name for kw in ["兩個內容", "two content"]):
                if "two_column" not in layout_map:
                    layout_map["two_column"] = idx

        # Fallbacks
        if "title_cover" not in layout_map:
            layout_map["title_cover"] = 0
        if "title_content" not in layout_map:
            layout_map["title_content"] = min(1, len(layouts) - 1)
        if "section" not in layout_map:
            layout_map["section"] = min(2, len(layouts) - 1)
        if "title_only" not in layout_map:
            layout_map["title_only"] = layout_map["title_content"]
        if "blank" not in layout_map:
            layout_map["blank"] = min(len(layouts) - 1, 6)
        if "two_column" not in layout_map:
            layout_map["two_column"] = layout_map["title_content"]
        return layout_map

    def _get_smart_layout(self, content: SlideContent):
        lt = content.layout_type
        has_chart = content.chart and content.chart.categories and content.chart.data_series
        is_first = content.page_number == 1
        is_last = content.page_number == len(self._plan.slides)
        is_divider = (lt == "title_slide" and not is_first)

        if is_first or is_last:
            idx = self._layout_map["title_cover"]
        elif is_divider:
            idx = self._layout_map["section"]
        elif lt == "full_chart" and has_chart:
            idx = self._layout_map.get("title_only", self._layout_map["title_content"])
        elif has_chart:
            idx = self._layout_map.get("title_only", self._layout_map["title_content"])
        else:
            idx = self._layout_map["title_content"]

        layouts = self._prs.slide_layouts
        layout = layouts[idx] if idx < len(layouts) else layouts[0]
        return layout, idx

    def _create_slide(self, content: SlideContent):
        """Get or create slide for this content page.
        
        Template mode: slides already exist (prepared in __init__), just return them.
        Blank mode: create new slide with layout.
        """
        if self._has_template:
            # Template mode: slides already prepared, just return the correct one
            slide_idx = content.page_number - 1  # 0-indexed
            if slide_idx < len(self._prs.slides):
                slide = self._prs.slides[slide_idx]
                # Clear existing text content from placeholders (keep decorative shapes)
                self._clear_slide_text(slide)
                return slide
            else:
                # Shouldn't happen, but fallback to adding a new slide
                layout = self._prs.slide_layouts[0]
                return self._prs.slides.add_slide(layout)
        else:
            # Blank mode: create new slide from layout
            layout, _ = self._get_smart_layout(content)
            return self._prs.slides.add_slide(layout)

    def _clear_slide_text(self, slide):
        """Clear text from placeholders on a template slide, keeping background shapes."""
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Check if it's a placeholder (using is_placeholder property)
                try:
                    if shape.is_placeholder:
                        shape.text_frame.clear()
                except (AttributeError, ValueError):
                    pass  # Not a placeholder, skip

    def generate(self) -> bytes:
        logger.info(f"PptxGenerator v3.0 (template-first) - generating {len(self._plan.slides)} slides, template={self._has_template}")
        for sc in self._plan.slides:
            self._add_slide(sc)
        buf = BytesIO()
        self._prs.save(buf)
        buf.seek(0)
        logger.info("PPTX generation complete with speaker notes on all slides")
        return buf.read()

    def _add_slide(self, content: SlideContent):
        lt = content.layout_type
        if lt == "title_slide":
            if content.page_number == 1:
                self._add_title_slide(content)
            else:
                self._add_chapter_divider(content)
        elif lt == "full_chart":
            self._add_full_chart_slide(content)
        elif lt == "content_with_chart":
            self._add_content_with_chart_slide(content)
        elif lt == "two_column":
            self._add_two_column_slide(content)
        else:
            self._add_content_only_slide(content)

    def _add_speaker_notes(self, slide, content: SlideContent):
        """Add AI analysis reasoning to speaker notes using robust XML approach."""
        notes_parts = []

        if content.insight_driver:
            notes_parts.append(f"【AI 分析驅動因素】\n{content.insight_driver}")

        if content.bullet_points:
            notes_parts.append(f"【關鍵要點】\n" + "\n".join(f"• {bp}" for bp in content.bullet_points))

        if content.chart and content.chart.title:
            notes_parts.append(f"【圖表說明】\n圖表類型：{content.chart.chart_type}\n圖表標題：{content.chart.title}")
            if content.chart.categories:
                notes_parts.append(f"資料維度：{', '.join(content.chart.categories[:5])}{'...' if len(content.chart.categories) > 5 else ''}")

        notes_parts.append(f"【頁面資訊】\n第 {content.page_number} 頁 | 版面類型：{content.layout_type}")

        notes_text = "\n\n".join(notes_parts)

        # Method 1: Standard python-pptx approach
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            # Clear existing paragraphs first
            tf.clear()
            # Write line by line as separate paragraphs for better formatting
            lines = notes_text.split("\n")
            for i, line in enumerate(lines):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line
            logger.info(f"Speaker notes added to page {content.page_number} (standard method)")
            return
        except Exception as e:
            logger.warning(f"Standard notes method failed for page {content.page_number}: {e}")

        # Method 2: Direct XML manipulation fallback
        try:
            self._add_notes_via_xml(slide, notes_text, content.page_number)
            logger.info(f"Speaker notes added to page {content.page_number} (XML fallback)")
            return
        except Exception as e:
            logger.error(f"XML notes fallback also failed for page {content.page_number}: {e}")

    def _add_notes_via_xml(self, slide, notes_text: str, page_number: int):
        """Add speaker notes by directly manipulating the slide's XML structure."""
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.slide import NotesSlide

        # Check if notes slide already exists
        try:
            notes_slide = slide.notes_slide
        except Exception:
            # If notes_slide property fails, create one manually
            notes_slide = self._create_notes_slide_manual(slide)

        if notes_slide is None:
            logger.warning(f"Cannot create notes slide for page {page_number}")
            return

        # Find the notes text frame body element
        # Notes slide has a shape with type "body" (idx=1) for the text
        sp_tree = notes_slide._element.find(qn('p:cSld')).find(qn('p:spTree'))

        # Find the body placeholder (usually ph type="body" or idx="1")
        body_sp = None
        for sp in sp_tree.findall(qn('p:sp')):
            nvSpPr = sp.find(qn('p:nvSpPr'))
            if nvSpPr is not None:
                nvPr = nvSpPr.find(qn('p:nvPr'))
                if nvPr is not None:
                    ph = nvPr.find(qn('p:ph'))
                    if ph is not None:
                        ph_type = ph.get('type', '')
                        ph_idx = ph.get('idx', '0')
                        if ph_type == 'body' or ph_idx == '1':
                            body_sp = sp
                            break

        if body_sp is None:
            # Try finding any text body that's not the slide image placeholder
            for sp in sp_tree.findall(qn('p:sp')):
                txBody = sp.find(qn('p:txBody'))
                if txBody is not None:
                    nvSpPr = sp.find(qn('p:nvSpPr'))
                    if nvSpPr is not None:
                        nvPr = nvSpPr.find(qn('p:nvPr'))
                        if nvPr is not None:
                            ph = nvPr.find(qn('p:ph'))
                            if ph is not None and ph.get('type', '') != 'sldImg':
                                body_sp = sp
                                break

        if body_sp is None:
            logger.warning(f"No body placeholder found in notes slide for page {page_number}")
            return

        # Clear existing text body and write new content
        txBody = body_sp.find(qn('p:txBody'))
        if txBody is None:
            return

        # Remove all existing paragraphs
        for p in txBody.findall(qn('a:p')):
            txBody.remove(p)

        # Add new paragraphs from notes_text
        for line in notes_text.split("\n"):
            p_elem = etree.SubElement(txBody, qn('a:p'))
            r_elem = etree.SubElement(p_elem, qn('a:r'))
            # Add run properties with font size
            rPr = etree.SubElement(r_elem, qn('a:rPr'))
            rPr.set('lang', 'zh-TW')
            rPr.set('sz', '1000')  # 10pt
            t_elem = etree.SubElement(r_elem, qn('a:t'))
            t_elem.text = line

    def _create_notes_slide_manual(self, slide):
        """Attempt to create a notes slide manually if the standard approach fails."""
        try:
            # Force creation through the presentation part
            prs_part = slide.part.package.presentation_part
            notes_master = prs_part.presentation.find(qn('p:notesMasterIdLst'))
            if notes_master is None:
                logger.warning("No notes master in presentation - cannot create notes slide")
                return None
            # Try the standard approach one more time after ensuring master exists
            return slide.notes_slide
        except Exception as e:
            logger.warning(f"Manual notes slide creation failed: {e}")
            return None

    # === Slide Type Methods ===

    def _add_title_slide(self, content: SlideContent):
        """Cover page — title centered, avoid bottom area."""
        slide = self._create_slide(content)
        if slide.shapes.title:
            slide.shapes.title.text = content.title
        else:
            self._add_text_box(slide, content.title,
                               Inches(2), Inches(2.5), Inches(9), Inches(1.5), Pt(32), bold=True)
        if content.subtitle:
            self._add_text_box(slide, content.subtitle,
                               Inches(2), Inches(4.2), Inches(9), Inches(0.8), Pt(16))
        self._add_speaker_notes(slide, content)

    def _add_chapter_divider(self, content: SlideContent):
        """Section divider — centered title on section layout."""
        slide = self._create_slide(content)
        if slide.shapes.title:
            slide.shapes.title.text = content.title
        else:
            self._add_text_box(slide, content.title,
                               Inches(3), Inches(3), Inches(7), Inches(1.5), Pt(28), bold=True)
        self._add_speaker_notes(slide, content)

    def _add_content_with_chart_slide(self, content: SlideContent):
        """Left: bullet points, Right: chart. Both within safe zone."""
        slide = self._create_slide(content)
        self._set_slide_title(slide, content.title)

        # Left: bullets (within safe zone)
        if content.bullet_points:
            bullets = "\n\n".join(f"• {bp}" for bp in content.bullet_points[:5])
            self._add_text_box(slide, bullets,
                               LEFT_MARGIN, CONTENT_TOP,
                               Inches(5.0), CONTENT_HEIGHT, Pt(11))

        # Right: chart (within safe zone, enough space for legend)
        if content.chart:
            self._add_native_chart(slide, content.chart,
                                   Inches(6.0), CONTENT_TOP,
                                   Inches(6.7), CONTENT_HEIGHT)

        self._add_speaker_notes(slide, content)

    def _add_full_chart_slide(self, content: SlideContent):
        """Full-width chart within safe content zone."""
        slide = self._create_slide(content)
        self._set_slide_title(slide, content.title)

        # Chart: full width but within safe vertical zone
        if content.chart:
            self._add_native_chart(slide, content.chart,
                                   LEFT_MARGIN, CONTENT_TOP,
                                   FULL_CONTENT_WIDTH, Inches(3.8))

        # Insight text below chart (still in safe zone)
        if content.insight_driver:
            self._add_text_box(slide, f"◆ {content.insight_driver}",
                               LEFT_MARGIN, Inches(5.4),
                               FULL_CONTENT_WIDTH, Inches(0.4), Pt(9), italic=True)

        self._add_speaker_notes(slide, content)

    def _add_two_column_slide(self, content: SlideContent):
        """Two-column: bullets left, chart/table right."""
        slide = self._create_slide(content)
        self._set_slide_title(slide, content.title)

        # Left column
        if content.bullet_points:
            bullets = "\n\n".join(f"• {bp}" for bp in content.bullet_points[:5])
            self._add_text_box(slide, bullets,
                               LEFT_MARGIN, CONTENT_TOP,
                               Inches(5.8), CONTENT_HEIGHT, Pt(11))

        # Right column
        if content.chart:
            self._add_native_chart(slide, content.chart,
                                   Inches(6.8), CONTENT_TOP,
                                   Inches(5.9), CONTENT_HEIGHT)
        elif content.table:
            self._add_table(slide, content.table.headers, content.table.rows,
                           left=Inches(6.8), top=CONTENT_TOP,
                           width=Inches(5.9), height=CONTENT_HEIGHT)

        self._add_speaker_notes(slide, content)

    def _add_content_only_slide(self, content: SlideContent):
        """Text-only slide within safe zone."""
        slide = self._create_slide(content)
        self._set_slide_title(slide, content.title)

        # Bullets in safe zone
        if content.bullet_points:
            bullets = "\n\n".join(f"• {bp}" for bp in content.bullet_points)
            self._add_text_box(slide, bullets,
                               LEFT_MARGIN, CONTENT_TOP,
                               FULL_CONTENT_WIDTH, Inches(3.6), Pt(12))

        # Insight in safe zone
        if content.insight_driver:
            self._add_text_box(slide, f"◆ 驅動因素：{content.insight_driver}",
                               LEFT_MARGIN, Inches(5.2),
                               FULL_CONTENT_WIDTH, Inches(0.5), Pt(10), italic=True)

        self._add_speaker_notes(slide, content)

    # === Chart & Table ===

    def _add_native_chart(self, slide, chart_spec: ChartSpec, left, top, width, height):
        """Add native editable chart with improved readability."""
        if not chart_spec.categories or not chart_spec.data_series:
            return
        try:
            ct = PPTX_CHART_TYPE_MAP.get(chart_spec.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            cd = CategoryChartData()
            # Limit categories for readability (max 12)
            categories = chart_spec.categories[:12]
            cd.categories = categories

            for s in chart_spec.data_series:
                v = s.values[:len(categories)]
                while len(v) < len(categories):
                    v.append(0)
                cd.add_series(s.name, tuple(v))

            cf = slide.shapes.add_chart(ct, left, top, width, height, cd)
            ch = cf.chart

            # Legend: always show for pie charts, show for multiple series otherwise
            is_pie = chart_spec.chart_type == "pie"
            if is_pie or len(chart_spec.data_series) > 1:
                ch.has_legend = True
                from pptx.enum.chart import XL_LEGEND_POSITION
                ch.legend.position = XL_LEGEND_POSITION.BOTTOM
                ch.legend.include_in_layout = False
                ch.legend.font.size = Pt(8)
            else:
                ch.has_legend = False

            # Chart title: concise
            if chart_spec.title:
                ch.has_title = True
                ch.chart_title.text_frame.text = chart_spec.title[:40]
                ch.chart_title.text_frame.paragraphs[0].font.size = Pt(9)
                ch.chart_title.text_frame.paragraphs[0].font.bold = True

            # Improve axis readability (not applicable for pie)
            if not is_pie:
                try:
                    if hasattr(ch, 'category_axis'):
                        ch.category_axis.tick_labels.font.size = Pt(8)
                    if hasattr(ch, 'value_axis'):
                        ch.value_axis.tick_labels.font.size = Pt(8)
                        ch.value_axis.has_major_gridlines = True
                except Exception:
                    pass

            logger.info(f"Chart added: '{chart_spec.title}' ({chart_spec.chart_type})")
        except Exception as e:
            logger.warning(f"Chart failed '{chart_spec.title}': {e}")

    def _add_table(self, slide, headers, rows, **pos):
        """Add table with clean formatting."""
        rc = min(len(rows) + 1, 12)  # Max 12 rows for readability
        cc = len(headers)
        t = slide.shapes.add_table(rc, cc, pos["left"], pos["top"], pos["width"], pos["height"]).table
        for i, h in enumerate(headers):
            t.cell(0, i).text = h
            t.cell(0, i).text_frame.paragraphs[0].font.bold = True
            t.cell(0, i).text_frame.paragraphs[0].font.size = Pt(9)
        for ri, row in enumerate(rows[:rc - 1], 1):
            for ci, val in enumerate(row):
                if ci < cc:
                    t.cell(ri, ci).text = str(val)
                    t.cell(ri, ci).text_frame.paragraphs[0].font.size = Pt(8)

    # === Utilities ===

    def _set_slide_title(self, slide, title: str):
        """Set title and clear empty placeholders to avoid overlap."""
        # Use the existing title placeholder
        if slide.shapes.title:
            slide.shapes.title.text = title
            for para in slide.shapes.title.text_frame.paragraphs:
                para.font.size = Pt(22)
                para.font.bold = True
        else:
            self._add_text_box(slide, title,
                               LEFT_MARGIN, TITLE_TOP,
                               FULL_CONTENT_WIDTH, TITLE_HEIGHT, Pt(22), bold=True)

        # Clear empty placeholders (set text to empty, make font tiny/transparent)
        for shape in list(slide.placeholders):
            if shape == slide.shapes.title:
                continue
            tf_text = shape.text_frame.text.strip() if shape.has_text_frame else ""
            if tf_text == "" or "按一下" in tf_text or "Click" in tf_text or "新增" in tf_text:
                # Clear the placeholder text completely
                shape.text_frame.clear()
                # Make it minimal size so it doesn't interfere
                shape.width = Inches(0.1)
                shape.height = Inches(0.1)

    def _add_text_box(self, slide, text, left, top, width, height,
                      font_size=Pt(12), bold=False, italic=False):
        """Add text box with proper formatting and overflow protection."""
        # Truncate text to prevent overflow (estimate based on area)
        max_chars = int((width / Inches(1)) * (height / Inches(1)) * 25)
        if len(text) > max_chars:
            text = text[:max_chars - 3] + "..."

        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True

        # Split by newlines and add as separate paragraphs for proper spacing
        lines = text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = font_size
            p.font.bold = bold
            p.font.italic = italic
            p.space_after = Pt(4)
